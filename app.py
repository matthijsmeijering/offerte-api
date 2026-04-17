from flask import Flask, request, send_file, jsonify
from pptx import Presentation
import requests
import os
import io

app = Flask(__name__)

TEMPLATE_URL = os.environ.get("TEMPLATE_URL", "")

def replace_in_paragraph(para, replacements):
    full_text = "".join(run.text for run in para.runs)
    needs_replace = any("{{" + k + "}}" in full_text for k in replacements)
    if not needs_replace:
        return
    new_text = full_text
    for key, value in replacements.items():
        placeholder = "{{" + key + "}}"
        new_text = new_text.replace(placeholder, str(value) if value else "")
    if para.runs:
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""

@app.route("/health", methods=["GET"])
def health():
    return jsonify({"status": "ok"})

@app.route("/generate", methods=["POST"])
def generate():
    data = request.get_json()
    if not data:
        return jsonify({"error": "Geen data meegestuurd"}), 400
try:
        session = requests.Session()
        r = session.get(TEMPLATE_URL, timeout=15)
        for key, value in r.cookies.items():
            if key.startswith('download_warning'):
                params = {'confirm': value, 'id': TEMPLATE_URL.split('id=')[-1]}
                r = session.get('https://drive.google.com/uc', params=params, timeout=15)
                break
        r.raise_for_status()
        template_bytes = io.BytesIO(r.content)
    except Exception as e:
        return jsonify({"error": f"Template ophalen mislukt: {str(e)}"}), 500
    prs = Presentation(template_bytes)
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    replace_in_paragraph(para, data)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        for para in cell.text_frame.paragraphs:
                            replace_in_paragraph(para, data)
    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    klantnaam = data.get("KLANTNAAM", "Klant").replace(" ", "_")
    datum = data.get("DATUM", "").replace(" ", "-")
    filename = f"Offerte_{klantnaam}_{datum}.pptx"
    return send_file(
        output,
        mimetype="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        as_attachment=True,
        download_name=filename
    )

if __name__ == "__main__":
    app.run(host="0.0.0.0", port=10000)
